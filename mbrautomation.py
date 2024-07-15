
import boto3
import pandas as pd
import io
from datetime import datetime, timedelta, timezone
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication

# Initialize AWS clients
ec2 = boto3.client('ec2')
s3 = boto3.client('s3')
iam = boto3.client('iam')
rds = boto3.client('rds')
cloudwatch = boto3.client('cloudwatch')
acm = boto3.client('acm')
ce_client = boto3.client('ce')
savingsplans = boto3.client('savingsplans')
config_client = boto3.client('config')
guardduty_client = boto3.client('guardduty')
securityhub_client = boto3.client('securityhub')
guardduty = boto3.client('guardduty')
inspector = boto3.client('inspector2')
config = boto3.client('config')

def get_users_without_access_key():
    iam = boto3.client('iam')
    response = iam.list_users()
    users_without_access_key = []

    for user in response['Users']:
        user_name = user['UserName']
        # Check if the user has any access keys
        access_keys = iam.list_access_keys(UserName=user_name)
        if not access_keys['AccessKeyMetadata']:
            users_without_access_key.append(user_name)

    return users_without_access_key

def is_mfa_enabled(user_name):
    try:
        iam = boto3.client('iam')
        mfa_status = iam.list_mfa_devices(UserName=user_name)
        if not mfa_status['MFADevices']:
            return 'disabled'
        return 'enabled'  # Just in case you need to distinguish enabled users
    except Exception as e:
        print(f"Error retrieving MFA status for user {user_name}: {e}")
        return 'unknown'

def get_mfa_enabled_console_users():
    users_without_access_key = get_users_without_access_key()

    if not users_without_access_key:
        # If no users without access keys are found, return a DataFrame with 'N/A' values
        return pd.DataFrame(columns=['User', 'MFA_STATUS'], data=[['N/A', 'N/A']])

    users_without_mfa = [(user, is_mfa_enabled(user)) for user in users_without_access_key]

    # Create a DataFrame
    df = pd.DataFrame(users_without_mfa, columns=['User', 'MFA_STATUS'])
    return df

def get_users_without_access_key():
    iam = boto3.client('iam')
    users = []

    paginator = iam.get_paginator('list_users')
    for response in paginator.paginate():
        for user in response['Users']:
            users.append(user['UserName'])

    return users

def password_creation_age(username):
    iam = boto3.client('iam')
    response = iam.get_user(UserName=username)
    user_create_date = response['User']['CreateDate']
    if user_create_date:
        return (datetime.now() - user_create_date.replace(tzinfo=None)).days
    else:
        return None

def access_key_age(username):
    iam = boto3.client('iam')
    response = iam.list_access_keys(UserName=username)
    if response['AccessKeyMetadata']:
        access_key_create_date = response['AccessKeyMetadata'][0]['CreateDate']
        return (datetime.now() - access_key_create_date.replace(tzinfo=None)).days
    else:
        return None

import pandas as pd

def get_inactive_users(threshold_days):
    inactive_users = []
    # Calculate the threshold date
    threshold_date = datetime.now(timezone.utc) - timedelta(days=threshold_days)
    print("Threshold Date:", threshold_date)

    # Get all IAM users
    response = iam.list_users()
    print("Number of Users Fetched:", len(response['Users']))

    # Iterate over each IAM user
    for user in response['Users']:
        user_name = user['UserName']
        user_inactive = True

        # Check the last login date for the user
        try:
            password_last_used = iam.get_user(UserName=user_name)['User'].get('PasswordLastUsed')
            if password_last_used and password_last_used >= threshold_date:
                user_inactive = False
        except KeyError:
            # If the 'PasswordLastUsed' key is not present, consider the user as inactive
            pass

        # Get a list of access keys for the user
        access_keys = iam.list_access_keys(UserName=user_name)['AccessKeyMetadata']

        for key in access_keys:
            access_key_id = key['AccessKeyId']
            last_used_response = iam.get_access_key_last_used(AccessKeyId=access_key_id)
            last_used_time = last_used_response['AccessKeyLastUsed'].get('LastUsedDate')

            if last_used_time and last_used_time >= threshold_date:
                user_inactive = False
                break

        # If both password and access keys are older than threshold days, mark the user as inactive
        if user_inactive:
            inactive_users.append({'User Name': user_name})

    # Create DataFrame from the inactive user data
    if inactive_users:
        return pd.DataFrame(inactive_users)
    else:
        # If no inactive users are found, return a DataFrame with 'N/A' values in every column
        return pd.DataFrame(columns=['User Name'], data=[['N/A']])

def get_security_group_details():
    try:
        response = ec2.describe_security_groups()
        sg_data = []
        for sg in response['SecurityGroups']:
            for rule in sg['IpPermissions']:
                if 'FromPort' in rule and rule['FromPort'] not in [80, 443]:
                    is_all_traffic = False
                    source = ""
                    if 'IpRanges' in rule:
                        for ip_range in rule['IpRanges']:
                            if ip_range['CidrIp'] == '0.0.0.0/0' or ip_range['CidrIp'] == '::/0':
                                is_all_traffic = True
                                source = "All Traffic"
                                break
                    if 'UserIdGroupPairs' in rule:
                        for pair in rule['UserIdGroupPairs']:
                            if pair['GroupId'] == '0.0.0.0/0' or pair['GroupId'] == '::/0':
                                is_all_traffic = True
                                source = "All Traffic"
                                break
                    if is_all_traffic:
                        sg_name = sg['GroupName']
                        sg_id = sg['GroupId']
                        port = rule['FromPort']
                        port_text = 'all port' if port in [0, -1] else str(port)
                        sg_data.append([sg_name, sg_id, port_text, source])
        if sg_data:
            return pd.DataFrame(sg_data, columns=['Security Group Name', 'Security Group ID', 'Port', 'Source'])
        else:
            # If no security group details are found, return a DataFrame with 'N/A' values in every column
            return pd.DataFrame(columns=['Security Group Name', 'Security Group ID', 'Port', 'Source'], data=[['N/A', 'N/A', 'N/A', 'N/A']])
    except Exception as e:
        print(f"An error occurred in get_security_group_details: {e}")
        raise e



def get_available_volumes():
    try:
        response = ec2.describe_volumes(
            Filters=[{'Name': 'status', 'Values': ['available']}]
        )
        volumes = response['Volumes']
        volume_data = [[vol['VolumeId'], vol['Size'], vol['State']] for vol in volumes]
        if volume_data:
            return pd.DataFrame(volume_data, columns=['Volume ID', 'Size (GB)', 'State'])
        else:
            # If no volumes are found, return a DataFrame with 'N/A' values in every column
            return pd.DataFrame(columns=['Volume ID', 'Size (GB)', 'State'], data=[['N/A', 'N/A', 'N/A']])
    except Exception as e:
        print(f"An error occurred in get_available_volumes: {e}")
        raise e


def get_unencrypted_volumes():
    try:
        response = ec2.describe_volumes()
        unencrypted_volumes = []

        for volume in response.get('Volumes', []):
            if not volume.get('Encrypted'):
                unencrypted_volumes.append([
                    volume.get('VolumeId', 'N/A'),
                    volume.get('Size', 'N/A'),
                    volume.get('VolumeType', 'N/A'),
                    'Unencrypted'
                ])

        if unencrypted_volumes:
            return pd.DataFrame(unencrypted_volumes, columns=['Volume ID', 'Size (GiB)', 'Volume Type', 'Encryption Status'])
        else:
            # If no unencrypted volumes are found, return a DataFrame with 'N/A' values in every column
            return pd.DataFrame(columns=['Volume ID', 'Size (GiB)', 'Volume Type', 'Encryption Status'],
                                data=[['N/A', 'N/A', 'N/A', 'N/A']])
    except Exception as e:
        print(f"An error occurred in get_unencrypted_volumes: {e}")
        raise e

def fetch_unassociated_ips():
    try:
        ec2 = boto3.client('ec2')
        response = ec2.describe_addresses()
        addresses = response.get('Addresses', [])

        ip_data = [
            {"Allocation ID": addr.get("AllocationId", "N/A"), "Public IP": addr.get("PublicIp", "N/A")}
            for addr in addresses
            if "AssociationId" not in addr
        ]

        if ip_data:
            unassociated_ips_df = pd.DataFrame(ip_data)
        else:
            # If no unassociated IP addresses are found, return a DataFrame with 'N/A' values in every column
            unassociated_ips_df = pd.DataFrame(columns=['Allocation ID', 'Public IP'], data=[['N/A', 'N/A']])

        # Reset the index to get the expected 0-based index
        unassociated_ips_df.reset_index(drop=True, inplace=True)

        return unassociated_ips_df
    except Exception as e:
        print(f"An error occurred in fetch_unassociated_ips: {e}")
        raise e

def get_ec2_instance_utilization():
    try:
        ec2 = boto3.client('ec2')
        cloudwatch = boto3.client('cloudwatch')
        utilization_data = []
        reservations = ec2.describe_instances()['Reservations']

        for reservation in reservations:
            for instance in reservation['Instances']:
                instance_id = instance['InstanceId']
                status_check = ec2.describe_instance_status(InstanceIds=[instance_id])
                if status_check['InstanceStatuses']:
                    instance_status = status_check['InstanceStatuses'][0]['InstanceStatus']['Status']
                    system_status = status_check['InstanceStatuses'][0]['SystemStatus']['Status']
                else:
                    instance_status = 'Unknown'
                    system_status = 'Unknown'

                # Query CPU utilization metric data
                cpu_metric_data = cloudwatch.get_metric_data(
                    MetricDataQueries=[
                        {
                            'Id': 'cpu_utilization',
                            'MetricStat': {
                                'Metric': {
                                    'Namespace': 'AWS/EC2',
                                    'MetricName': 'CPUUtilization',
                                    'Dimensions': [
                                        {'Name': 'InstanceId', 'Value': instance_id}
                                    ]
                                },
                                'Period': 3600,
                                'Stat': 'Average',
                                'Unit': 'Percent'
                            },
                            'ReturnData': True
                        }
                    ],
                    StartTime=datetime.utcnow() - timedelta(days=1),
                    EndTime=datetime.utcnow()
                )
                cpu_utilization = cpu_metric_data['MetricDataResults'][0]['Values']
                avg_cpu_utilization = sum(cpu_utilization) / len(cpu_utilization) if cpu_utilization else 'N/A'

                # Query disk utilization metric data
                disk_metric_data = cloudwatch.get_metric_data(
                    MetricDataQueries=[
                        {
                            'Id': 'disk_utilization',
                            'MetricStat': {
                                'Metric': {
                                    'Namespace': 'AWS/EC2',
                                    'MetricName': 'DiskSpaceUtilization',
                                    'Dimensions': [
                                        {'Name': 'InstanceId', 'Value': instance_id}
                                    ]
                                },
                                'Period': 3600,
                                'Stat': 'Average',
                                'Unit': 'Percent'
                            },
                            'ReturnData': True
                        }
                    ],
                    StartTime=datetime.utcnow() - timedelta(days=1),
                    EndTime=datetime.utcnow()
                )
                disk_utilization = disk_metric_data['MetricDataResults'][0]['Values']
                avg_disk_utilization = sum(disk_utilization) / len(disk_utilization) if disk_utilization else 'N/A'

                # Query memory utilization metric data
                memory_metric_data = cloudwatch.get_metric_data(
                    MetricDataQueries=[
                        {
                            'Id': 'memory_utilization',
                            'MetricStat': {
                                'Metric': {
                                    'Namespace': 'AWS/EC2',
                                    'MetricName': 'MemoryUtilization',
                                    'Dimensions': [
                                        {'Name': 'InstanceId', 'Value': instance_id}
                                    ]
                                },
                                'Period': 3600,
                                'Stat': 'Average',
                                'Unit': 'Percent'
                            },
                            'ReturnData': True
                        }
                    ],
                    StartTime=datetime.utcnow() - timedelta(days=1),
                    EndTime=datetime.utcnow()
                )
                memory_utilization = memory_metric_data['MetricDataResults'][0]['Values']
                avg_memory_utilization = sum(memory_utilization) / len(memory_utilization) if memory_utilization else 'N/A'

                # Add data to the list
                utilization_data.append([
                    instance_id, instance_status, system_status,
                    avg_cpu_utilization, avg_disk_utilization, avg_memory_utilization
                ])

        # Create a DataFrame and return it
        columns = ['Instance ID', 'Instance Status', 'System Status', 'Avg CPU Utilization (%)', 'Avg Disk Utilization (%)', 'Avg Memory Utilization(%)']
        return pd.DataFrame(utilization_data, columns=columns) if utilization_data else pd.DataFrame(columns=columns, data=[['N/A'] * len(columns)])
    except Exception as e:
        print(f"An error occurred in get_ec2_instance_utilization: {e}")
        raise e

# Function to check engine end-of-life status (Dummy function, replace with actual logic)
def check_engine_end_of_life(engine, engine_version):
    # Dummy implementation, replace with actual end-of-life checking logic
    return "No"


def get_rds_utilization(threshold=90):
    try:
        rds = boto3.client('rds')
        cloudwatch = boto3.client('cloudwatch')

        instances = rds.describe_db_instances()['DBInstances']

        utilization_data = []
        for instance in instances:
            instance_id = instance['DBInstanceIdentifier']

            # Query CPU utilization metric data
            metrics = cloudwatch.get_metric_data(
                MetricDataQueries=[
                    {
                        'Id': 'cpu_query',
                        'MetricStat': {
                            'Metric': {
                                'Namespace': 'AWS/RDS',
                                'MetricName': 'CPUUtilization',
                                'Dimensions': [
                                    {'Name': 'DBInstanceIdentifier', 'Value': instance_id}
                                ]
                            },
                            'Period': 3600,
                            'Stat': 'Average',
                            'Unit': 'Percent'
                        },
                        'ReturnData': True
                    },
                    {
                        'Id': 'freeable_memory_query',
                        'MetricStat': {
                            'Metric': {
                                'Namespace': 'AWS/RDS',
                                'MetricName': 'FreeableMemory',
                                'Dimensions': [
                                    {'Name': 'DBInstanceIdentifier', 'Value': instance_id}
                                ]
                            },
                            'Period': 3600,
                            'Stat': 'Average',
                            'Unit': 'Bytes'
                        },
                        'ReturnData': True
                    },
                    {
                        'Id': 'free_storage_space_query',
                        'MetricStat': {
                            'Metric': {
                                'Namespace': 'AWS/RDS',
                                'MetricName': 'FreeStorageSpace',
                                'Dimensions': [
                                    {'Name': 'DBInstanceIdentifier', 'Value': instance_id}
                                ]
                            },
                            'Period': 3600,
                            'Stat': 'Average',
                            'Unit': 'Bytes'
                        },
                        'ReturnData': True
                    }
                ],
                StartTime=datetime.utcnow() - timedelta(days=1),  # Adjust time period if needed
                EndTime=datetime.utcnow()
            )

            # Process metric data
            cpu_utilization = metrics['MetricDataResults'][0]['Values']
            avg_cpu_utilization = sum(cpu_utilization) / len(cpu_utilization) if cpu_utilization else 'N/A'
            freeable_memory = metrics['MetricDataResults'][1]['Values']
            avg_freeable_memory = sum(freeable_memory) / len(freeable_memory) if freeable_memory else 'N/A'
            free_storage_space = metrics['MetricDataResults'][2]['Values']
            avg_free_storage_space = sum(free_storage_space) / len(free_storage_space) if free_storage_space else 'N/A'

            # Convert memory and storage to GB
            avg_freeable_memory_gb = avg_freeable_memory / (1024 ** 3) if avg_freeable_memory != 'N/A' else 'N/A'
            avg_free_storage_space_gb = avg_free_storage_space / (1024 ** 3) if avg_free_storage_space != 'N/A' else 'N/A'

            utilization_data.append([instance_id, avg_cpu_utilization, avg_freeable_memory_gb, avg_free_storage_space_gb])

            # Check if CPU utilization exceeds the threshold
            if isinstance(avg_cpu_utilization, float) and avg_cpu_utilization > threshold:
                print(f"Instance {instance_id} has CPU utilization above {threshold}%.")

        columns = ['Instance ID', 'Average CPU Utilization (%)', 'Average Freeable Memory (GB)', 'Average Free Storage Space (GB)']
        return pd.DataFrame(utilization_data, columns=columns) if utilization_data else pd.DataFrame(columns=columns, data=[['N/A'] * len(columns)])
    except Exception as e:
        print(f"An error occurred in get_rds_utilization: {e}")
        raise e



# Function to get RDS engine versions and end-of-life status

def get_rds_engine_versions():
    try:
        rds = boto3.client('rds')

        instances = rds.describe_db_instances().get('DBInstances', [])

        engine_data = []
        for instance in instances:
            instance_id = instance.get('DBInstanceIdentifier', 'N/A')
            engine = instance.get('Engine', 'N/A')
            engine_version = instance.get('EngineVersion', 'N/A')

            # Check engine end-of-life status
            end_of_life = check_engine_end_of_life(engine, engine_version)

            # Fetch certificate details
            certificate_expiry = 'N/A'
            certificate_identifier = instance.get('CACertificateIdentifier')
            if certificate_identifier:
                try:
                    certificate_info = rds.describe_certificates(CertificateIdentifier=certificate_identifier)
                    certificate_expiry = certificate_info['Certificates'][0]['ValidTill'].strftime('%Y-%m-%d')
                except Exception as e:
                    print(f"Error fetching certificate info: {e}")

            # Get RDS status and storage auto-scaling status
            status = instance.get('DBInstanceStatus', 'N/A')
            auto_scaling_enabled = 'Enabled' if instance.get('StorageAutoscaling') else 'Disabled'

            engine_data.append([instance_id, engine, engine_version, end_of_life, certificate_expiry, status, auto_scaling_enabled])

        columns = ['Instance ID', 'Engine', 'Engine Version', 'End of Life', 'DB Instance Certificate Expiration Date', 'Status', 'Storage Auto-Scaling Enabled']
        return pd.DataFrame(engine_data, columns=columns) if engine_data else pd.DataFrame(columns=columns, data=[['N/A'] * len(columns)])
    except Exception as e:
        print(f"An error occurred in get_rds_engine_versions: {e}")
        raise e


def get_ec2_reservations():
    try:
        ec2 = boto3.client('ec2')

        reservations = ec2.describe_reserved_instances().get('ReservedInstances', [])

        reservations_data = []
        for reservation in reservations:
            if reservation['State'] != 'active':
                continue
            reservation_id = reservation.get('ReservedInstancesId', 'N/A')
            instance_type = reservation.get('InstanceType', 'N/A')
            end = reservation.get('End', 'N/A')
            instance_count = reservation.get('InstanceCount', 'N/A')
            payment_option = reservation.get('UsagePrice', 'N/A')
            offering_class = reservation.get('OfferingClass', 'N/A')
            scope = reservation.get('Scope', 'N/A')
            platform = reservation.get('ProductDescription', 'N/A')
            reservations_data.append([reservation_id, instance_type, end.strftime('%Y-%m-%d') if isinstance(end, datetime) else 'N/A', instance_count, payment_option, offering_class, scope, platform])

        columns = ['Reservation ID', 'Instance Type', 'Expiry Date', 'Instance Count', 'Payment Option', 'Offering Class', 'Scope', 'Platform']
        return pd.DataFrame(reservations_data, columns=columns) if reservations_data else pd.DataFrame(columns=columns, data=[['N/A'] * len(columns)])
    except Exception as e:
        print(f"An error occurred in get_ec2_reservations: {e}")
        raise e

def get_savings_plans():
    try:
        savings_plans = savingsplans.describe_savings_plans().get('savingsPlans', [])

        plans_data = []
        for plan in savings_plans:
            plan_id = plan.get('savingsPlanId', 'N/A')
            plan_type = plan.get('savingsPlanType', 'N/A')
            instance_family = plan.get('instanceFamily', 'N/A')
            region = plan.get('region', 'N/A')
            end = plan.get('end', 'N/A')
            plans_data.append([plan_id, plan_type, instance_family, region, end])

        columns = ['Plan ID', 'Plan Type', 'Instance Family', 'Region', 'End Date']
        return pd.DataFrame(plans_data, columns=columns) if plans_data else pd.DataFrame(columns=columns, data=[['N/A'] * len(columns)])
    except Exception as e:
        print(f"An error occurred in get_savings_plans: {e}")
        raise e


def get_rds_reservations():
    try:
        reservations = rds.describe_reserved_db_instances().get('ReservedDBInstances', [])

        if not reservations:
            return pd.DataFrame(columns=[
                'Reservation ID', 'Product', 'Region', 'Class', 'Remaining Days', 'Multi-AZ', 'Quantity', 'Offering Type'
            ], data=[['N/A'] * 8])

        reservations_data = []
        for reservation in reservations:
            if reservation['State'] != 'active':
                continue
            reservation_id = reservation.get('ReservedDBInstanceId', 'N/A')
            product = reservation.get('ProductDescription', 'N/A')
            region = reservation.get('AvailabilityZone', 'N/A')
            instance_class = reservation.get('DBInstanceClass', 'N/A')
            end_date = reservation.get('StartTime', datetime.utcnow()) + timedelta(seconds=reservation.get('Duration', 0))
            remaining_days = (end_date - datetime.utcnow()).days
            multi_az = reservation.get('MultiAZ', 'N/A')
            quantity = reservation.get('DBInstanceCount', 'N/A')
            offering_type = reservation.get('OfferingType', 'N/A')
            reservations_data.append([
                reservation_id, product, region, instance_class, remaining_days, multi_az, quantity, offering_type
            ])

        columns = ['Reservation ID', 'Product', 'Region', 'Class', 'Remaining Days', 'Multi-AZ', 'Quantity', 'Offering Type']
        return pd.DataFrame(reservations_data, columns=columns)
    except Exception as e:
        print(f"An error occurred in get_rds_reservations: {e}")
        raise e

def check_acm_expiry():
    # List all certificates in the account
    response = acm.list_certificates()

    # Track ACM data
    acm_data = []

    # Get details for each certificate
    for cert in response['CertificateSummaryList']:
        cert_details = acm.describe_certificate(CertificateArn=cert['CertificateArn'])

        # Check if 'Certificate' key exists in the response
        if 'Certificate' in cert_details:
            expiration_date = cert_details['Certificate'].get('NotAfter')

            # Check if expiration date exists
            if expiration_date:
                expiration_date = expiration_date.replace(tzinfo=None)  # Remove timezone info

                # Check if the certificate is expired or expiring soon (within 30 days)
                if expiration_date < datetime.utcnow():
                    status = "Expired"
                elif (expiration_date - datetime.utcnow()) <= timedelta(days=30):
                    status = "Expiring soon (within 30 days)"
                else:
                    status = "Valid"

                # Extract domain name
                domain_name = cert_details['Certificate'].get('DomainName')

                # Check if domain name exists
                if domain_name:
                    # Append certificate data
                    acm_data.append([cert_details['Certificate']['CertificateArn'], domain_name, expiration_date, status])
                else:
                    # Append certificate data with no domain name
                    acm_data.append([cert_details['Certificate']['CertificateArn'], "N/A", expiration_date, "No domain name"])
            else:
                # Append certificate data with no expiration date
                acm_data.append([cert_details['Certificate']['CertificateArn'], "N/A", "N/A", "No expiration date"])
        else:
            # Append certificate data with failed status domain verification
            acm_data.append([cert['CertificateArn'], "N/A", "N/A", "Failed domain verification"])

    # Check if acm_data is empty and append N/A values if necessary
    if not acm_data:
        acm_data.append(["N/A", "N/A", "N/A", "N/A"])

    # Create a DataFrame and return it
    return pd.DataFrame(acm_data, columns=['Certificate ARN', 'Domain Name', 'Expiration Date', 'Expiration Status'])


def check_budget_alarms():
    try:
        # Create a boto3 client for the AWS Budgets service
        budgets_client = boto3.client('budgets')

        # Retrieve the account ID
        account_id = boto3.client('sts').get_caller_identity()['Account']

        # Describe budgets for the account
        response = budgets_client.describe_budgets(AccountId=account_id)
        budgets = response.get('Budgets', [])

        if budgets:
            # If there are budgets, return a DataFrame with budget names and their status as 'Set'
            return pd.DataFrame([[budget['BudgetName'], 'Set'] for budget in budgets], columns=['Budget Name', 'Status'])
        else:
            # If no budgets are found, return a DataFrame indicating no budget alarms are set
            return pd.DataFrame([['N/A', 'N/A']], columns=['Budget Name', 'Status'])
    except Exception as e:
        print(f"An error occurred in check_budget_alarms: {e}")
        raise e


def get_ec2_termination_protection():
    try:
        # Retrieve the termination protection information for each EC2 instance
        response = ec2.describe_instances()
        instance_data = []
        for reservation in response.get('Reservations', []):
            for instance in reservation.get('Instances', []):
                instance_id = instance['InstanceId']
                instance_state = instance['State']['Name']
                termination_protection = instance.get('InstanceLifecycle') == 'spot' or instance.get('InstanceLifecycle') == 'scheduled'
                instance_data.append([instance_id, instance_state, 'Yes' if termination_protection else 'No'])
        if instance_data:
            return pd.DataFrame(instance_data, columns=['Instance ID', 'Instance State', 'Termination Protection'])
        else:
            return pd.DataFrame(columns=['Instance ID', 'Instance State', 'Termination Protection']).fillna('N/A')
    except Exception as e:
        print(f"An error occurred in get_ec2_termination_protection: {e}")
        raise e



def get_security_hub_status():
    try:
        response = securityhub_client.describe_hub()
        status = 'ACTIVE' if 'HubArn' in response else 'INACTIVE'
        security_hub_status = [[
            response.get('HubArn', 'N/A'),
            response.get('SubscribedAt', 'N/A'),
            status
        ]]
    except securityhub_client.exceptions.InvalidAccessException:
        security_hub_status = [['N/A', 'N/A', 'INACTIVE']]

    return pd.DataFrame(security_hub_status, columns=['Hub ARN', 'Subscribed At', 'Status'])

def get_config_status():
    config_status = []

    try:
        response = config_client.describe_configuration_recorders()
        if response['ConfigurationRecorders']:
            for recorder in response['ConfigurationRecorders']:
                recorder_name = recorder.get('name', 'N/A')

                status_response = config_client.describe_configuration_recorder_status(ConfigurationRecorderNames=[recorder_name])
                for status in status_response['ConfigurationRecordersStatus']:
                    recording = status.get('recording', 'N/A')
                    last_status = status.get('lastStatus', 'N/A')

                    config_status.append([
                        recorder_name,
                        'ACTIVE' if recording else 'INACTIVE',
                        last_status
                    ])
        else:
            config_status.append(['N/A', 'INACTIVE', 'N/A'])

    except Exception as e:
        config_status.append(['N/A', 'ERROR', str(e)])

    return pd.DataFrame(config_status, columns=['Recorder Name', 'Status', 'Last Status'])

def download_image_from_s3(bucket_name, key):
    s3 = boto3.client('s3')
    obj = s3.get_object(Bucket=bucket_name, Key=key)
    return io.BytesIO(obj['Body'].read())


def save_data_to_ppt(dataframes, function_names, bucket_name, logo_key, first_slide_image_key, last_slide_image_key, escalation_key):
    # Create a PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(13.2)
    prs.slide_height = Inches(7.5)
 
    # Fetch images from S3
    first_slide_image = download_image_from_s3(bucket_name, first_slide_image_key)
    logo_image = download_image_from_s3(bucket_name, logo_key)
    last_slide_image = download_image_from_s3(bucket_name, last_slide_image_key)
    escalation_image = download_image_from_s3(bucket_name, escalation_key)
 
    # Add the first slide with the specific image
    first_slide_layout = prs.slide_layouts[6]  # Use a blank layout
    first_slide = prs.slides.add_slide(first_slide_layout)
    
 
    # Add the specific image to the first slide
    first_slide_width = prs.slide_width
    first_slide_height = prs.slide_height
    first_slide.shapes.add_picture(first_slide_image, 0, 0, width=first_slide_width, height=first_slide_height)
 
    # Prompt for the input string
    name = "Equity Pandit"
 
    # Add the name to the first slide
    left = Inches(1)
    top = Inches(3)  # Move the name down
    width = Inches(8)
    height = Inches(1)
    text_box = first_slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # White font color
    p.font.name = 'Helvetica'
 
    # Add the date of execution below the name with less gap
    top = Inches(3.5)  # Move the date down accordingly
    text_box = first_slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    execution_date = datetime.now().strftime("%Y-%m-%d")
    p = text_frame.add_paragraph()
    p.text = execution_date
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # White font color
    p.font.name = 'Helvetica'
 
    # Add slides for each dataframe
    for dataframe, function_name in zip(dataframes, function_names):
        slide_layout = prs.slide_layouts[5]  # Use a layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        # Set the slide background color to light grey
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light grey background
        #Title details
        title = slide.shapes.title
        title.left = Inches(1)
        title.top = Inches(0.8)
        title.width =  Inches(11)
        title.text = function_name
        title.text_frame.paragraphs[0].alignment = 1 # Align text to the left
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Set font size for title to 24
        title.text_frame.paragraphs[0].font.name = 'Helvetica'  # Set font to Helvetica
        title.text_frame.paragraphs[0].font.bold = True  # Make the title bold
 
        # Add the logo to the slide
        logo_width = Inches(0.6)
        logo_height = Inches(0.6)
        left = Inches(0.4)  # Position logo 0.3 inch from the left edge
        top = Inches(6.7)  # Position logo 0.7 inch from the top edge
        slide.shapes.add_picture(logo_image, left, top, width=logo_width, height=logo_height)
 
        # Adjust this value based on your table width
        left = Inches(1)  
        top = Inches(1.2)
        width = Inches(11.2)
        height = Inches(6)
 
        try:
            table = slide.shapes.add_table(dataframe.shape[0] + 1, dataframe.shape[1], left, top, width, height).table
        except ValueError as e:
            print(f"Error adding table for {function_name}: {e}")
            continue

        # Set the height for the first row
        first_row_height = Inches(0.7)  # Set the desired height for the first row
        table.rows[0].height = first_row_height

        # Set font size to 14 and font to Helvetica for the table text
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(13)
                    run.font.name = 'Helvetica'
        
        # Set first row background color to  blue 69,  and font color to white
        cobalt_blue = RGBColor(0, 70, 170)
        for col_index, column_name in enumerate(dataframe.columns):
            cell = table.cell(0, col_index)
            cell.text = column_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = cobalt_blue
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(15)
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White font color

        # Add data to the table
        for row_index, row in enumerate(dataframe.itertuples(), start=1):
            for col_index, value in enumerate(row[1:], start=0):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(13)
                        run.font.name = 'Helvetica'


    # Add the last slide with the specific image
    last_slide_layout = prs.slide_layouts[6]  # Use a blank layout
    escalation_slide_layout = prs.slide_layouts[6]    #escalation layout
    escalation_slide = prs.slides.add_slide(escalation_slide_layout)
    last_slide = prs.slides.add_slide(last_slide_layout)

    # Add the specific image to the last slide
    last_slide_width = escalation_slide_width = prs.slide_width
    last_slide_height = escalation_slide_height = prs.slide_height
    last_slide.shapes.add_picture(last_slide_image, 0, 0, width=first_slide_width, height=first_slide_height)
    escalation_slide.shapes.add_picture(escalation_image, 0, 0, width=first_slide_width, height=first_slide_height)

    # Save the presentation to a buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
 
    # Upload PowerPoint file to S3
    s3 = boto3.client('s3')
    s3.put_object(Bucket=bucket_name, Key='UIchanges_aws_data.pptx', Body=ppt_buffer)
    print('data.pptx has been saved to S3 bucket "{}".'.format(bucket_name))

def send_email_with_attachment(bucket_name, object_key, sender, recipient):
    ses = boto3.client('ses', region_name='us-east-1')  # Replace with your SES region

    # Fetch the file from S3
    s3 = boto3.client('s3')
    response = s3.get_object(Bucket=bucket_name, Key=object_key)
    ppt_data = response['Body'].read()

    # Create the email message
    msg = MIMEMultipart()
    msg['Subject'] = 'AWS Data PowerPoint Presentation'
    msg['From'] = sender
    msg['To'] = recipient

    body = MIMEText('Please find the attached PowerPoint presentation with the AWS data.', 'plain')
    msg.attach(body)

    attachment = MIMEApplication(ppt_data)
    attachment.add_header('Content-Disposition', 'attachment', filename=object_key)
    msg.attach(attachment)

    # Send the email using SES
    ses.send_raw_email(
        Source=sender,
        Destinations=[recipient],
        RawMessage={'Data': msg.as_string()}
    )

def lambda_handler():
    threshold_days = 90

    # Get data from different functions
    dataframes = [
        get_mfa_enabled_console_users(),
        get_inactive_users(threshold_days=threshold_days),
        get_security_group_details(),
        get_available_volumes(),
        get_unencrypted_volumes(),
        get_ec2_termination_protection(),
        fetch_unassociated_ips(),
        get_ec2_instance_utilization(),
        get_rds_engine_versions(),
        get_rds_utilization(),
        get_ec2_reservations(),
        get_savings_plans(),
        get_rds_reservations(),
        check_acm_expiry(),
        check_budget_alarms(),
        get_security_hub_status(),
        get_config_status(),
    ]

    # S3 bucket and object keys for the images
    bucket_name = 'govardhanicici'  # Update with your bucket name
    logo_key = 'c4c.jpg'  # Update with the actual key to the logo in S3
    first_slide_image_key = 'c4c2.jpg'  # Update with the actual key to the first slide image in S3
    last_slide_image_key = 'thankyou.jpg'  # Update with the actual key to the first slide image in S3
    escalation_key = 'escalation.jpg'   #Update escalation matrix

    # Save combined data to a PowerPoint presentation
    save_data_to_ppt(dataframes, [
        'Console users without MFA',
        'Inactive Users',
        'Security groups',
        'Available Volumes',
        'Unencrypted Volumes',
        'EC2 Termination Protection',
        'Unassociated IP',
        'EC2 Instance Utilization',
        'RDS Details',
        'RDS Utilization',
        'EC2 Reservations',
        'Saving Plans',
        'RDS Reservations',
        'ACM',
        'Budget alarms',
        'Security Hub',
        'Config',
    ], bucket_name, logo_key, first_slide_image_key, last_slide_image_key, escalation_key)

    # Retrieve sender and recipient from environment variables with defaults
    sender_email = os.getenv('SENDER_EMAIL', 'default-sender@example.com')
    receiver_email = os.getenv('RECEIVER_EMAIL', 'default-receiver@example.com')

    send_email_with_attachment(
        bucket_name=bucket_name,
        object_key='UIchanges_aws_data.pptx',
        sender=sender_email,
        recipient=receiver_email
    )

    return {
        'statusCode': 200,
        'body': 'Data has been saved to an AWS PowerPoint presentation.'
    }

lambda_handler()
